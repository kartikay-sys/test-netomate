import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation (16:9 aspect ratio)
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# --- Theme Colors ---
BG_DARK = RGBColor(15, 23, 42)      # Slate 900
BOX_DARK = RGBColor(30, 41, 59)     # Slate 800
TEXT_WHITE = RGBColor(248, 250, 252)
TEXT_GRAY = RGBColor(148, 163, 184)
ACCENT_BLUE = RGBColor(56, 189, 248)
ACCENT_GREEN = RGBColor(52, 211, 153)
ACCENT_PURPLE = RGBColor(167, 139, 250)

def set_background(slide):
    """Draw a dark rectangle over the entire slide as background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return bg

def add_title(slide, text, subtitle=None):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(40)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = TEXT_GRAY
        p2.alignment = PP_ALIGN.LEFT

def create_model_box(slide, left, top, width, height, title, cost_in, cost_out, req_1k, req_10k, req_100k, mo_200, mo_1k, mo_5k, accent_color, highlight=False):
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = BOX_DARK
    box.line.color.rgb = accent_color if highlight else RGBColor(51, 65, 85)
    if highlight:
        box.line.width = Pt(3)

    # Content
    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = accent_color
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = f"Cost/1M Tok: In: ${cost_in:.4f} | Out: ${cost_out:.4f}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY
    
    # Bulk Section
    p = tf.add_paragraph()
    p.text = "\nBulk Requests:"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    
    tf.add_paragraph().text = f"• 1k requests: ${req_1k}"
    tf.add_paragraph().text = f"• 10k requests: ${req_10k}"
    tf.add_paragraph().text = f"• 100k requests: ${req_100k}"
    
    # Monthly Section
    p = tf.add_paragraph()
    p.text = "\nMonthly Projections (30d):"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    
    tf.add_paragraph().text = f"• 200/day: ${mo_200}"
    tf.add_paragraph().text = f"• 1k/day: ${mo_1k}"
    tf.add_paragraph().text = f"• 5k/day: ${mo_5k}"
    
    for paragraph in tf.paragraphs[3:]:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = TEXT_GRAY
        if highlight and "1k" in paragraph.text or "5k" in paragraph.text:
            pass


# ----------------------------------------------------
# Slide 1: Title
# ----------------------------------------------------
layout_blank = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(layout_blank)
set_background(slide1)

txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "NetoMate AI Architecture"
p.font.size = Pt(54)
p.font.color.rgb = TEXT_WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Cost Analysis & Optimal LLM Selection for Telecom Test Flow Orchestration"
p2.font.size = Pt(24)
p2.font.color.rgb = TEXT_GRAY
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = "\nBaseline Context: 4,121 In | 994 Out"
p3.font.size = Pt(18)
p3.font.color.rgb = ACCENT_BLUE
p3.alignment = PP_ALIGN.CENTER

# ----------------------------------------------------
# Slide 2: Gemini
# ----------------------------------------------------
slide2 = prs.slides.add_slide(layout_blank)
set_background(slide2)
add_title(slide2, "Google Gemini Models", "Monthly 30-Day Projections (USD)")

create_model_box(slide2, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "Gemini 2.5 Flash", 0.0012, 0.0025, 
                 "3.72", "37.21", "372.07", "22.32", "111.62", "558.11", ACCENT_GREEN)

create_model_box(slide2, Inches(4.75), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "Gemini 3 Flash", 0.0021, 0.0030, 
                 "5.04", "50.42", "504.18", "30.25", "151.25", "756.27", ACCENT_BLUE)

create_model_box(slide2, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "Gemini 3 Pro", 0.0082, 0.0119, 
                 "20.17", "201.67", "2,016.71", "121.00", "605.01", "3,025.07", ACCENT_PURPLE)

# ----------------------------------------------------
# Slide 3: OpenAI
# ----------------------------------------------------
slide3 = prs.slides.add_slide(layout_blank)
set_background(slide3)
add_title(slide3, "OpenAI Models", "Highlighting Cost Efficiency at Scale")

create_model_box(slide3, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "GPT-4o Mini (Budget)", 0.0006, 0.0006, 
                 "1.21", "12.15", "121.45", "7.29", "36.44", "182.18", ACCENT_GREEN, highlight=True)

create_model_box(slide3, Inches(4.75), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "GPT-5.4 Mini", 0.0031, 0.0045, 
                 "7.56", "75.64", "756.38", "45.38", "226.91", "1,134.56", ACCENT_BLUE)

create_model_box(slide3, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "GPT-5.4 (Flagship)", 0.0103, 0.0149, 
                 "25.21", "252.12", "2,521.25", "151.27", "756.38", "3,781.88", ACCENT_PURPLE)

# ----------------------------------------------------
# Slide 4: Anthropic
# ----------------------------------------------------
slide4 = prs.slides.add_slide(layout_blank)
set_background(slide4)
add_title(slide4, "Anthropic Next-Gen Models", "Monthly 30-Day Projections (USD)")

create_model_box(slide4, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "Haiku 4.5", 0.0041, 0.0050, 
                 "9.09", "90.91", "909.10", "54.55", "272.73", "1,363.65", ACCENT_GREEN)

create_model_box(slide4, Inches(4.75), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "Sonnet 4.6", 0.0124, 0.0149, 
                 "27.27", "272.73", "2,727.30", "163.64", "818.19", "4,090.95", ACCENT_BLUE)

create_model_box(slide4, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.5), 
                 "Opus 4.7", 0.0206, 0.0249, 
                 "45.46", "454.55", "4,545.50", "272.73", "1,363.65", "6,818.25", ACCENT_PURPLE)


# ----------------------------------------------------
# Slide 5: Recommendation
# ----------------------------------------------------
slide5 = prs.slides.add_slide(layout_blank)
set_background(slide5)
add_title(slide5, "Optimal Model Selection", "Why GPT-4o Mini is the perfect fit for NetoMate's Architecture")

box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
box.fill.solid()
box.fill.fore_color.rgb = BOX_DARK
box.line.color.rgb = ACCENT_GREEN
box.line.width = Pt(2)

txBox = slide5.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.33), Inches(4.4))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "The Clear Winner: GPT-4o Mini\n"
p.font.size = Pt(32)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "1. Native Structured Outputs (JSON):"
p2.font.bold = True
p2.font.size = Pt(20)
p2.font.color.rgb = TEXT_WHITE

p3 = tf.add_paragraph()
p3.text = "   Step D of your pipeline demands strict JSON without conversational filler. GPT-4o Mini natively supports 'response_format: json_object' yielding a 100% parse success rate.\n"
p3.font.size = Pt(18)
p3.font.color.rgb = TEXT_GRAY

p4 = tf.add_paragraph()
p4.text = "2. Leveraging Existing Guardrails:"
p4.font.bold = True
p4.font.size = Pt(20)
p4.font.color.rgb = TEXT_WHITE

p5 = tf.add_paragraph()
p5.text = "   Because NetoMate has rigid pre-processing (Step B state checks) and post-processing (Step E hallucination filters), you don't need a flagship model's reasoning. You just need a robust pattern matcher.\n"
p5.font.size = Pt(18)
p5.font.color.rgb = TEXT_GRAY

p6 = tf.add_paragraph()
p6.text = "3. Unbeatable Cost Efficiency:"
p6.font.bold = True
p6.font.size = Pt(20)
p6.font.color.rgb = TEXT_WHITE

p7 = tf.add_paragraph()
p7.text = "   At just $182/month for massive enterprise volume (5,000 req/day), it is effectively 7.5x cheaper than Haiku 4.5 and 15x cheaper than Gemini 3 Pro, leaving massive budget for UI/server infrastructure."
p7.font.size = Pt(18)
p7.font.color.rgb = TEXT_GRAY


# Save Presentation
output_path = r"e:\NETOMATE\NetoMate_Cost_Analysis.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
